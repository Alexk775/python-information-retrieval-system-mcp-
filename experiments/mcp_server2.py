#!/usr/bin/env python3
# mcp_server.py — V6, Robust Reranking

import os, time, hashlib, io, re, sys
from pathlib import Path
from typing import Dict, List, Optional

import fitz  # PDF parser (PyMuPDF)
from pptx import Presentation
from docx import Document
from PIL import Image
import pytesseract

from rapidfuzz import fuzz
# --- FIX 1: REMOVED problematic top-level imports ---
# These were causing the adapter to crash on import.
# from helpers.fuzzy_engine import smart_fuzzy_score, mmr_rerank
# from helpers.text_utils import normalize_text, token_estimate

# ------------------ CONFIG ------------------
DATA_DIR = Path(os.getenv("MCP_DATA_RAW", "data/raw"))
USE_OCR = os.getenv("MCP_USE_OCR", "1") != "0"
SUPPORTED_EXTS = {".pdf", ".pptx", ".docx", ".png", ".jpg", ".jpeg", ".txt", ".md"}

# --- Semantic knobs (relaxed & richer snippets) ---
SEM_ENABLED = os.getenv("MCP_SEMANTIC", "1") != "0"
SEM_ALPHA = float(os.getenv("MCP_SEM_ALPHA", "0.5"))
# --- V6: Widen the net to find more candidates ---
SEM_TOPK = int(os.getenv("MCP_SEM_TOPK", "100")) # Was 50
SEM_SNIPPET = int(os.getenv("MCP_SNIPPET", "600"))

# Fallback triggers
MIN_RESULTS = int(os.getenv("MCP_MIN_RESULTS", "3"))
MIN_TOP_SCORE = float(os.getenv("MCP_MIN_TOP_SCORE", "0.35"))
# --- V6: Removed LEXICAL_GUARANTEE_SCORE (it was unreliable) ---
MIN_SEM_BASE = float(os.getenv("MCP_MIN_SEM", "0.05"))

# Band, rescue, tail rules
SEM_MARGIN = float(os.getenv("MCP_SEM_MARGIN", "0.09"))
KEEP_MIN = int(os.getenv("MCP_KEEP_MIN", "2"))
KEEP_MAX = int(os.getenv("MCP_KEEP_MAX", "5"))  # tuned for top_k=4
RESCUE_SEM_BASE = float(os.getenv("MCP_RESCUE_SEM", "0.10"))
RESCUE_MAX_BASE = int(os.getenv("MCP_RESCUE_MAX", "1"))  # max 1 rescue
TAIL_KEEP = int(os.getenv("MCP_TAIL_KEEP", "1"))

# Generic demotion (auto-disabled when few docs exist)
GENERIC_PENALTY = float(os.getenv("MCP_GENERIC_PENALTY", "0.06"))
GENERIC_TITLES = {
    "company overview",
    "vision & mission",
    "vision and mission",
    "overview",
    "about us",
}

# Top-k (agent-facing)
TOP_K_DEFAULT = int(os.getenv("MCP_TOP_K", "4"))


def now_iso() -> str:
    return time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())


# ------------------ MODELS ------------------
class Resource:
    def __init__(
        self,
        uri: str,
        title: str,
        mime_type: str,
        path: Path,
        size_bytes: int,
        updated_at: str,
        collection: str,
        version: int = 1,
    ):
        self.uri = uri
        self.title = title
        self.mime_type = mime_type
        self.path = path
        self.size_bytes = size_bytes
        self.updated_at = updated_at
        self.collection = collection
        self.version = version

    def asdict(self) -> Dict:
        return {
            "uri": self.uri,
            "title": self.title,
            "mime_type": self.mime_type,
            "collections": [self.collection],
            "labels": [],
            "lang": "auto",
            "version": self.version,
            "size_bytes": self.size_bytes,
            "updated_at": self.updated_at,
            "path": str(self.path),
        }


class Index:
    def __init__(self):
        # resource base-uri -> Resource
        self.resources: Dict[str, Resource] = {}
        # full-uri (with locator) -> text
        self.chunks: Dict[str, str] = {}

    # ---------- Public API used by adapters ----------
    def list_resources(self, collection: Optional[str] = None) -> List[Resource]:
        vals = list(self.resources.values())
        if collection:
            vals = [r for r in vals if r.collection == collection]
        return sorted(vals, key=lambda r: r.title.lower())

    def read_resource_text(self, uri: str, offset: int = 0, limit: int = 4000):
        if uri not in self.chunks:
            # Try to find chunk by prefix if a base_uri was passed
            for chunk_uri in self.chunks.keys():
                if chunk_uri.startswith(uri):
                    uri = chunk_uri
                    break
            else:
                 raise KeyError(f"resource chunk not found: {uri}")

        content = self.chunks[uri]
        total = len(content)
        start = max(0, min(offset, total))
        end = min(total, start + limit)
        sliced = content[start:end]
        base_uri, locator = uri.split("#", 1) if "#" in uri else (uri, "")
        res = self.resources.get(base_uri)
        citations = [
            {
                "source": res.title if res else base_uri,
                "locator": locator or "full",
                "confidence": 0.9,
            }
        ]
        paging = {"offset": start, "limit": end - start, "total": total}
        return sliced, citations, paging

    def _need_semantic(self, rows: List[Dict]) -> bool:
        if not rows or len(rows) < MIN_RESULTS:
            return True
        try:
            top = max(r.get("score", 0.0) for r in rows)
        except ValueError:
            return True
        return top < MIN_TOP_SCORE

    def search(self, query: str, top_k: int = None, filters: dict | None = None):
        """Hybrid search: fuzzy+TFIDF (lexical) → optional semantic rerank → MMR."""
        # --- FIX 1: MOVED IMPORTS HERE ---
        # This "lazy loading" prevents the adapter from crashing on import.
        from helpers.auto_synonym_engine import AutoSynonymEngine
        from helpers.fuzzy_engine import smart_fuzzy_score, mmr_rerank

        if top_k is None:
            top_k = TOP_K_DEFAULT

        filters = filters or {}
        coll = set(filters.get("collections", []))

        # ---- adaptive relax based on corpus size ----
        total_docs = len(self.chunks)
        few_docs = total_docs < 12
        MIN_SEM = max(0.0, MIN_SEM_BASE - (0.03 if few_docs else 0.0))
        RESCUE_MAX = RESCUE_MAX_BASE + (1 if few_docs else 0)
        apply_generic_penalty = not few_docs

        # --- semantic expansion (TF-IDF/WordNet) ---
        engine = AutoSynonymEngine()
        corpus_texts = list(self.chunks.values())
        if corpus_texts:
            engine.fit(corpus_texts)
        expanded_terms = engine.expand_query(query, max_terms=8)
        expanded_query = " ".join(expanded_terms)

        # --- lexical scoring (fuzzy + tiny TF-IDF pairwise) ---
        rows: List[Dict] = []
        for uri, chunk in self.chunks.items():
            if coll:
                base = uri.split("#", 1)[0]
                meta = self.resources.get(base)
                if not meta or meta.collection not in coll:
                    continue

            # **NOTE**: We search the *full chunk* now, not just a snippet
            fuzzy_score = smart_fuzzy_score(expanded_query, chunk)
            # --- V6: Tighter lexical cutoff to reduce junk ---
            if fuzzy_score < 35:  # Was 20
                continue

            tfidf_score = 0
            if engine.vec is not None:
                from sklearn.feature_extraction.text import TfidfVectorizer

                vec = TfidfVectorizer(stop_words="english")
                vec.fit([expanded_query, chunk])
                tfidf_score = (
                    vec.transform([expanded_query])
                    .dot(vec.transform([chunk]).T)
                    .toarray()[0][0]
                    * 100
                )

            final_score = 0.7 * fuzzy_score + 0.3 * tfidf_score  # (0..100)
            lex_norm = final_score / 100.0
            rows.append(
                {
                    "uri": uri,
                    # Snippet is now based on the *start* of the small chunk
                    "snippet": chunk[:SEM_SNIPPET],
                    "score": round(lex_norm, 3),
                }
            )

        # --- sort + pre-cut for semantic ---
        rows.sort(key=lambda r: r["score"], reverse=True)
        rows = rows[:SEM_TOPK]

        # --- V6: Reverted to "Always Rerank" logic ---
        # The V5 "Trust" logic failed because the lexical score is unreliable.
        # We MUST run the semantic reranker to find non-keyword matches.
        used_semantic = False
        if SEM_ENABLED:
            try:
                # 1. Rerank the *entire* candidate pile
                from helpers.semantic_rerank import rerank_semantic
                sem_rows = rerank_semantic(
                    query, rows, top_k=max(top_k, KEEP_MIN), alpha=SEM_ALPHA
                )
                used_semantic = True

                # 2. Apply the semantic filters (floor, band, etc.)
                # 1) semantic floor
                sem_rows = [
                    r
                    for r in sem_rows
                    if r.get("semantic_score", 0.0) >= MIN_SEM
                ]

                band_rows: List[Dict] = []
                if sem_rows:
                    # 2) band around the best
                    max_sem = max(r["semantic_score"] for r in sem_rows)
                    band_cut = max_sem - SEM_MARGIN
                    band_rows = [
                        r for r in sem_rows if r["semantic_score"] >= band_cut
                    ]

                    # backfill if below KEEP_MIN
                    if len(band_rows) < KEEP_MIN:
                        sem_rows_sorted = sorted(
                            sem_rows,
                            key=lambda x: x["semantic_score"],
                            reverse=True,
                        )
                        band_rows = sem_rows_sorted[:KEEP_MIN]

                    # trim if above KEEP_MAX
                    if len(band_rows) > KEEP_MAX:
                        band_rows = sorted(
                            band_rows,
                            key=lambda x: x["semantic_score"],
                            reverse=True,
                        )[:KEEP_MAX]

                    # 3) Domain-aware light demotion for generic titles
                    if apply_generic_penalty and band_rows:
                        adjusted = []
                        for r in band_rows:
                            base = r["uri"].split("#", 1)[0]
                            title = (
                                (
                                    self.resources.get(base).title
                                    if self.resources.get(base)
                                    else ""
                                )
                                .lower()
                            )
                            is_generic = title in GENERIC_TITLES
                            if is_generic and r["semantic_score"] < (max_sem - 0.02):
                                r = {
                                    **r,
                                    "score": max(0.0, r["score"] - GENERIC_PENALTY),
                                }
                            adjusted.append(r)
                        band_rows = sorted(
                            adjusted, key=lambda x: x["score"], reverse=True
                        )

                    # 4) Rescue pass
                    if RESCUE_MAX > 0:
                        chosen_uris = {r["uri"] for r in band_rows}
                        rescue = []
                        for r in sem_rows:
                            if r["uri"] in chosen_uris:
                                continue
                            if r.get("semantic_score", 0.0) >= RESCUE_SEM_BASE:
                                rescue.append(r)
                            if len(rescue) >= RESCUE_MAX:
                                break
                        if rescue:
                            band_rows = sorted(
                                band_rows + rescue,
                                key=lambda x: x["semantic_score"],
                                reverse=True,
                            )

                    # 5) Tail keep
                    if len(band_rows) < KEEP_MIN and TAIL_KEEP > 0 and sem_rows:
                        chosen_uris = {r["uri"] for r in band_rows}
                        tail_candidates = [
                            r for r in sem_rows if r["uri"] not in chosen_uris
                        ]
                        if tail_candidates:
                            tail_best = max(
                                tail_candidates,
                                key=lambda x: (
                                    x.get("semantic_score", 0.0),
                                    x.get("score", 0.0),
                                ),
                            )
                            band_rows.append(tail_best)

                # 3. Final list is the filtered band
                if band_rows:
                    rows = band_rows
                else:
                    # Fallback to lexical if semantic filters cleared everything
                    rows = rows[:top_k]
                    for r in rows:
                        r["rerank"] = "lexical_fallback_v6"

            except Exception as e:
                print(f"Semantic reranking failed: {e}", file=sys.stderr)
                # Fallback to just the lexical rows if reranking blows up
                rows = rows[:top_k]
        else:
            # Semantic is disabled, just use lexical
            rows = rows[:top_k]

        # --- final MMR diversity ---
        # Run diversity *after* merging all our results
        rows = mmr_rerank(rows, lam=0.65, top_k=top_k)

        # normalize to 0..1 and tag rerank origin
        for r in rows:
            r["score"] = round(r["score"], 3)
            r["rerank"] = "semantic" if used_semantic else r.get("rerank", "lexical")

        return rows

    # ---------- Indexing ----------
    def index_all(self):
        self.resources.clear()
        self.chunks.clear()
        DATA_DIR.mkdir(parents=True, exist_ok=True)

        for path in sorted(DATA_DIR.rglob("*")):
            if not path.is_file():
                continue
            ext = path.suffix.lower()
            if ext not in SUPPORTED_EXTS:
                continue
            self._index_file(path)

    def _index_file(self, path: Path):
        ext = path.suffix.lower()
        collection = path.parent.name or "default"
        stat = path.stat()
        doc_id = hashlib.sha1(str(path).encode()).hexdigest()[:10]
        base_uri = f"mcp://{collection}/{doc_id}"
        title = path.stem
        updated_at = now_iso()

        res = Resource(
            uri=base_uri,
            title=title,
            mime_type=ext,
            path=path,
            size_bytes=stat.st_size,
            updated_at=updated_at,
            collection=collection,
            version=1,
        )
        self.resources[base_uri] = res

        if ext == ".pdf":
            self._parse_pdf(path, base_uri)
        elif ext == ".pptx":
            self._parse_pptx(path, base_uri)
        elif ext == ".docx":
            self._parse_docx(path, base_uri)
        elif ext in {".png", ".jpg", ".jpeg"}:
            self._parse_image(path, base_uri)
        elif ext in {".txt", ".md"}:
            self._parse_text(path, base_uri)

    # ---------- NEW CHUNKING HELPER ----------

    def _add_chunk(self, base_uri: str, locator: str, text: str):
        """
        A helper to normalize and add a text chunk to the index.
        Skips empty or very short chunks.
        """
        text = (text or "").strip()
        # Normalize whitespace (replace multiple newlines/tabs with a single space)
        # Keep single newlines for basic structure
        text = re.sub(r"[\t\r\f\v]+", " ", text)
        text = re.sub(r" +", " ", text)
        text = re.sub(r"\n\s*\n", "\n", text) # Collapse multiple newlines
        
        # Skip chunks that are too short to be meaningful
        if not text or len(text) < 20: # 20 char minimum
            return
        
        self.chunks[f"{base_uri}#{locator}"] = text

    # ---------- V3 PARSERS (Smarter Chunking) ----------

    def _parse_pdf(self, path: Path, base_uri: str):
        doc = fitz.open(path)
        try:
            for i, page in enumerate(doc, start=1):
                page_text = page.get_text("text") or ""
                page_text = (page_text or "").strip()

                if not page_text and USE_OCR:
                    pix = page.get_pixmap(dpi=300)
                    img = Image.open(io.BytesIO(pix.tobytes("png"))).convert("L")
                    try:
                        page_text = pytesseract.image_to_string(
                            img, lang="eng", config="--oem 1 --psm 6"
                        )
                    except Exception:
                        page_text = ""
                
                page_text = (page_text or "").strip()
                if not page_text:
                    continue  # skip empty pages

                # *** NEW CHUNKING ***
                # Split the page text by paragraphs (double newline)
                chunks = re.split(r"\n\s*\n", page_text)
                for j, chunk in enumerate(chunks):
                    self._add_chunk(base_uri, f"page-{i}-p-{j}", chunk)
        finally:
            doc.close()

    def _parse_pptx(self, path: Path, base_uri: str):
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides, start=1):
            chunk_idx = 0
            
            # Add each shape's text as a separate chunk
            for sp_idx, sp in enumerate(slide.shapes):
                if hasattr(sp, "text") and sp.text:
                    self._add_chunk(base_uri, f"slide-{i}-shape-{sp_idx}", sp.text)
                    chunk_idx += 1
                
                # Add table content as chunks (one per row)
                if hasattr(sp, "has_table") and sp.has_table:
                    try:
                        for row_idx, row in enumerate(sp.table.rows):
                            cells = [(cell.text or "").strip() for cell in row.cells]
                            line = " | ".join([c for c in cells if c])
                            self._add_chunk(base_uri, f"slide-{i}-table-{sp_idx}-row-{row_idx}", line)
                            chunk_idx += 1
                    except Exception:
                        pass

            # Add notes as their own chunk
            try:
                if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text or ""
                    self._add_chunk(base_uri, f"slide-{i}-notes", notes)
            except Exception:
                pass

    def _parse_docx(self, path: Path, base_uri: str):
        doc = Document(path)
        current_heading = "section-0"
        chunk_idx = 0

        # Add tables first, each row as a chunk
        try:
            for i, table in enumerate(doc.tables):
                for j, row in enumerate(table.rows):
                    cells = [(cell.text or "").strip() for cell in row.cells]
                    line = " | ".join([c for c in cells if c])
                    self._add_chunk(base_uri, f"table-{i}-row-{j}", line)
        except Exception:
            pass # Ignore table errors

        # Process paragraphs: each paragraph is a chunk
        for p in doc.paragraphs:
            txt = (p.text or "").strip()
            
            # If it's a heading, update our locator "namespace"
            if p.style and getattr(p.style, "name", "").startswith("Heading"):
                # Sanitize heading text to create a locator
                safe_heading = re.sub(r'[^a-zA-Z0-9]+', '-', txt.lower()[:30])
                current_heading = safe_heading or f"section-{chunk_idx}"
                chunk_idx = 0 # Reset chunk index for new section
                # Add the heading itself as a chunk
                self._add_chunk(base_uri, f"h-{current_heading}", txt)
            else:
                # Add the paragraph as a chunk
                self._add_chunk(base_uri, f"{current_heading}-p-{chunk_idx}", txt)
                chunk_idx += 1
                
        # Add headers/footers (as one chunk, they are usually small)
        try:
            hf_text = []
            for s in doc.sections:
                if s.header:
                    hf_text.extend([(p.text or "").strip() for p in s.header.paragraphs if (p.text or "").strip()])
                if s.footer:
                    hf_text.extend([(p.text or "").strip() for p in s.footer.paragraphs if (p.text or "").strip()])
            self._add_chunk(base_uri, "header-footer", "\n".join(hf_text))
        except Exception:
            pass


    def _parse_text(self, path: Path, base_uri: str):
        try:
            txt = path.read_text(encoding="utf-8", errors="ignore")
        except Exception:
            txt = path.read_text(errors="ignore")
        
        # *** NEW CHUNKING ***
        # Split the text by paragraphs (double newline)
        chunks = re.split(r"\n\s*\n", txt)
        for j, chunk in enumerate(chunks):
            self._add_chunk(base_uri, f"p-{j}", chunk)

    def _parse_image(self, path: Path, base_uri: str):
        text = ""
        if USE_OCR:
            try:
                img = Image.open(path).convert("L")
                text = pytesseract.image_to_string(
                    img, lang="eng", config="--oem 1 --psm 6"
                )
            except Exception:
                text = ""
        
        # Add the full OCR text as one chunk
        self._add_chunk(base_uri, "image-ocr", text)


# ------------ build index at import ------------
INDEX = Index()
INDEX.index_all()

# --- FIX 2: This now prints to stderr ---
# This is safe and will not break the JSON-RPC stdout pipe.
# This confirms the index is built *before* the adapter starts.
print(f"--- Index built. {len(INDEX.chunks)} chunks loaded. ---", file=sys.stderr)



