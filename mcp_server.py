#!/usr/bin/env python3
# mcp_server.py — tuned for mixed docs (PDF slides, DOCX, PPTX, TXT/MD)
import os, time, hashlib, io, re
from pathlib import Path
from typing import Dict, List, Optional

import fitz  # PDF parser (PyMuPDF)
from pptx import Presentation
from docx import Document
from PIL import Image
import pytesseract

from rapidfuzz import fuzz
from helpers.fuzzy_engine import smart_fuzzy_score, mmr_rerank
from helpers.text_utils import normalize_text, token_estimate

# ------------------ CONFIG ------------------
DATA_DIR = Path(os.getenv("MCP_DATA_RAW", "data/raw"))
USE_OCR = os.getenv("MCP_USE_OCR", "1") != "0"
SUPPORTED_EXTS = {".pdf", ".pptx", ".docx", ".png", ".jpg", ".jpeg", ".txt", ".md"}

# --- Semantic knobs (relaxed & richer snippets) ---
SEM_ENABLED = os.getenv("MCP_SEMANTIC", "1") != "0"
SEM_ALPHA = float(os.getenv("MCP_SEM_ALPHA", "0.5"))
### FIX 1: Widen the candidate pool for the reranker ###
SEM_TOPK = int(os.getenv("MCP_SEM_TOPK", "100")) # Was 50
SEM_SNIPPET = int(os.getenv("MCP_SEM_SNIPPET", "1500"))

# Fallback triggers
MIN_RESULTS = int(os.getenv("MCP_MIN_RESULTS", "3"))
MIN_TOP_SCORE = float(os.getenv("MCP_MIN_TOP_SCORE", "0.55")) # Was 0.35
MIN_SEM_BASE = float(os.getenv("MCP_MIN_SEM", "0.05"))

# Band, rescue, tail rules
SEM_MARGIN = float(os.getenv("MCP_SEM_MARGIN", "0.09"))
KEEP_MIN = int(os.getenv("MCP_KEEP_MIN", "2"))
KEEP_MAX = int(os.getenv("MCP_KEEP_MAX", "5"))  # tuned for top_k=4
RESCUE_SEM_BASE = float(os.getenv("MCP_RESCUE_SEM", "0.10"))
RESCUE_MAX_BASE = int(os.getenv("MCP_RESCUE_MAX", "1"))  # max 1 rescue
TAIL_KEEP = int(os.getenv("MCP_TAIL_KEEP", "1"))

# Generic demotion (auto-disabled when few docs exist)
GENERIC_PENALTY = float(os.getenv("MCP_GENERIC_PENALTY", "0.06"))
GENERIC_TITLES = {
    "company overview",
    "vision & mission",
    "vision and mission",
    "overview",
    "about us",
}

# Top-k (agent-facing)
TOP_K_DEFAULT = int(os.getenv("MCP_TOP_K", "4"))


def now_iso() -> str:
    return time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())


# ------------------ MODELS ------------------
class Resource:
    def __init__(
        self,
        uri: str,
        title: str,
        mime_type: str,
        path: Path,
        size_bytes: int,
        updated_at: str,
        collection: str,
        version: int = 1,
    ):
        self.uri = uri
        self.title = title
        self.mime_type = mime_type
        self.path = path
        self.size_bytes = size_bytes
        self.updated_at = updated_at
        self.collection = collection
        self.version = version

    def asdict(self) -> Dict:
        return {
            "uri": self.uri,
            "title": self.title,
            "mime_type": self.mime_type,
            "collections": [self.collection],
            "labels": [],
            "lang": "auto",
            "version": self.version,
            "size_bytes": self.size_bytes,
            "updated_at": self.updated_at,
            "path": str(self.path),
        }


class Index:
    def __init__(self):
        # resource base-uri -> Resource
        self.resources: Dict[str, Resource] = {}
        # full-uri (with locator) -> text
        self.chunks: Dict[str, str] = {}

    # ---------- Public API used by adapters ----------
    def list_resources(self, collection: Optional[str] = None) -> List[Resource]:
        vals = list(self.resources.values())
        if collection:
            vals = [r for r in vals if r.collection == collection]
        return sorted(vals, key=lambda r: r.title.lower())

    def read_resource_text(self, uri: str, offset: int = 0, limit: int = 4000):
        if uri not in self.chunks:
            raise KeyError(uri)
        content = self.chunks[uri]
        total = len(content)
        start = max(0, min(offset, total))
        end = min(total, start + limit)
        sliced = content[start:end]
        base_uri, locator = uri.split("#", 1) if "#" in uri else (uri, "")
        res = self.resources.get(base_uri)
        citations = [
            {
                "source": res.title if res else base_uri,
                "locator": locator or "full",
                "confidence": 0.9,
            }
        ]
        paging = {"offset": start, "limit": end - start, "total": total}
        return sliced, citations, paging

    def _need_semantic(self, rows: List[Dict]) -> bool:
        if not rows or len(rows) < MIN_RESULTS:
            return True
        try:
            top = max(r.get("score", 0.0) for r in rows)
        except ValueError:
            return True
        return top < MIN_TOP_SCORE

    def search(self, query: str, top_k: int = None, filters: dict | None = None):
        """Hybrid search: fuzzy+TFIDF (lexical) → optional semantic rerank → MMR."""
        from helpers.auto_synonym_engine import AutoSynonymEngine
        from helpers.fuzzy_engine import smart_fuzzy_score, mmr_rerank

        if top_k is None:
            top_k = TOP_K_DEFAULT

        filters = filters or {}
        coll = set(filters.get("collections", []))

        # ---- adaptive relax based on corpus size ----
        total_docs = len(self.chunks)
        few_docs = total_docs < 12
        MIN_SEM = max(0.0, MIN_SEM_BASE - (0.03 if few_docs else 0.0))
        RESCUE_MAX = RESCUE_MAX_BASE + (1 if few_docs else 0)
        apply_generic_penalty = not few_docs

        # --- semantic expansion (TF-IDF/WordNet) ---
        engine = AutoSynonymEngine()
        corpus_texts = list(self.chunks.values())
        if corpus_texts:
            engine.fit(corpus_texts)
        expanded_terms = engine.expand_query(query, max_terms=8)
        expanded_query = " ".join(expanded_terms)

        # --- lexical scoring (fuzzy + tiny TF-IDF pairwise) ---
        rows: List[Dict] = []
        for uri, chunk in self.chunks.items():
            if coll:
                base = uri.split("#", 1)[0]
                meta = self.resources.get(base)
                if not meta or meta.collection not in coll:
                    continue

            snippet_full = chunk[:4000]
            fuzzy_score = smart_fuzzy_score(expanded_query, snippet_full)
            # relaxed cutoff (was 40)
            if fuzzy_score < 20:
                continue

            tfidf_score = 0
            if engine.vec is not None:
                from sklearn.feature_extraction.text import TfidfVectorizer

                vec = TfidfVectorizer(stop_words="english")
                vec.fit([expanded_query, snippet_full])
                tfidf_score = (
                    vec.transform([expanded_query])
                    .dot(vec.transform([snippet_full]).T)
                    .toarray()[0][0]
                    * 100
                )

            final_score = 0.7 * fuzzy_score + 0.3 * tfidf_score  # (0..100)
            lex_norm = final_score / 100.0
            rows.append(
                {
                    "uri": uri,
                    "snippet": snippet_full[:SEM_SNIPPET],
                    "score": round(lex_norm, 3),
                }
            )

        # --- sort + pre-cut for semantic ---
        rows.sort(key=lambda r: r["score"], reverse=True)
        rows = rows[:SEM_TOPK]

        # --- optional semantic rerank / band / rescue / tail ---
        used_semantic = False
        if SEM_ENABLED and self._need_semantic(rows):
            try:
                from helpers.semantic_rerank import rerank_semantic

                sem_rows = rerank_semantic(
                    query, rows, top_k=max(top_k, KEEP_MIN), alpha=SEM_ALPHA
                )

                # 1) semantic floor
                sem_rows = [
                    r
                    for r in sem_rows
                    if r.get("semantic_score", 0.0) >= MIN_SEM
                ]

                band_rows: List[Dict] = []
                if sem_rows:
                    # 2) band around the best
                    max_sem = max(r["semantic_score"] for r in sem_rows)
                    band_cut = max_sem - SEM_MARGIN
                    band_rows = [
                        r for r in sem_rows if r["semantic_score"] >= band_cut
                    ]

                    # backfill if below KEEP_MIN
                    if len(band_rows) < KEEP_MIN:
                        sem_rows_sorted = sorted(
                            sem_rows,
                            key=lambda x: x["semantic_score"],
                            reverse=True,
                        )
                        band_rows = sem_rows_sorted[:KEEP_MIN]

                    # trim if above KEEP_MAX
                    if len(band_rows) > KEEP_MAX:
                        band_rows = sorted(
                            band_rows,
                            key=lambda x: x["semantic_score"],
                            reverse=True,
                        )[:KEEP_MAX]

                    # 3) Domain-aware light demotion for generic titles
                    if apply_generic_penalty and band_rows:
                        adjusted = []
                        for r in band_rows:
                            base = r["uri"].split("#", 1)[0]
                            title = (
                                (
                                    self.resources.get(base).title
                                    if self.resources.get(base)
                                    else ""
                                )
                                .lower()
                            )
                            is_generic = title in GENERIC_TITLES
                            if is_generic and r["semantic_score"] < (
                                max_sem - 0.02
                            ):
                                r = {
                                    **r,
                                    "score": max(
                                        0.0, r["score"] - GENERIC_PENALTY
                                    ),
                                }
                            adjusted.append(r)
                        band_rows = sorted(
                            adjusted, key=lambda x: x["score"], reverse=True
                        )

                    # 4) Rescue pass
                    if RESCUE_MAX > 0:
                        chosen_uris = {r["uri"] for r in band_rows}
                        rescue = []
                        for r in sem_rows:
                            if r["uri"] in chosen_uris:
                                continue
                            if r.get("semantic_score", 0.0) >= RESCUE_SEM_BASE:
                                rescue.append(r)
                            if len(rescue) >= RESCUE_MAX:
                                break
                        if rescue:
                            band_rows = sorted(
                                band_rows + rescue,
                                key=lambda x: x["semantic_score"],
                                reverse=True,
                            )

                    # 5) Tail keep
                    if (
                        len(band_rows) < KEEP_MIN
                        and TAIL_KEEP > 0
                        and sem_rows
                    ):
                        chosen_uris = {r["uri"] for r in band_rows}
                        tail_candidates = [
                            r for r in sem_rows if r["uri"] not in chosen_uris
                        ]
                        if tail_candidates:
                            tail_best = max(
                                tail_candidates,
                                key=lambda x: (
                                    x.get("semantic_score", 0.0),
                                    x.get("score", 0.0),
                                ),
                            )
                            band_rows.append(tail_best)

                if sem_rows and band_rows:
                    rows = band_rows
                    used_semantic = True
                else:
                    rows = rows[:top_k]
                    for r in rows:
                        r["rerank"] = "lexical_fallback"
            except Exception:
                rows = rows[:top_k]
        else:
            rows = rows[:top_k]

        # --- final MMR diversity ---
        rows = mmr_rerank(rows, lam=0.65, top_k=top_k)

        # normalize to 0..1 and tag rerank origin
        for r in rows:
            r["score"] = round(r["score"], 3)
            r["rerank"] = "semantic" if used_semantic else r.get("rerank", "lexical")

        return rows

    # ---------- Indexing ----------
    def index_all(self):
        self.resources.clear()
        self.chunks.clear()
        DATA_DIR.mkdir(parents=True, exist_ok=True)

        for path in sorted(DATA_DIR.rglob("*")):
            if not path.is_file():
                continue
            ext = path.suffix.lower()
            if ext not in SUPPORTED_EXTS:
                continue
            self._index_file(path)

    def _index_file(self, path: Path):
        ext = path.suffix.lower()
        collection = path.parent.name or "default"
        stat = path.stat()
        doc_id = hashlib.sha1(str(path).encode()).hexdigest()[:10]
        base_uri = f"mcp://{collection}/{doc_id}"
        title = path.stem
        updated_at = now_iso()

        res = Resource(
            uri=base_uri,
            title=title,
            mime_type=ext,
            path=path,
            size_bytes=stat.st_size,
            updated_at=updated_at,
            collection=collection,
            version=1,
        )
        self.resources[base_uri] = res

        if ext == ".pdf":
            self._parse_pdf(path, base_uri)
        elif ext == ".pptx":
            self._parse_pptx(path, base_uri)
        elif ext == ".docx":
            self._parse_docx(path, base_uri)
        elif ext in {".png", ".jpg", ".jpeg"}:
            self._parse_image(path, base_uri)
        elif ext in {".txt", ".md"}:
            self._parse_text(path, base_uri)

    # ---------- Parsers ----------
    def _parse_pdf(self, path: Path, base_uri: str):
        doc = fitz.open(path)
        try:
            for i, page in enumerate(doc, start=1):
                # 1) Try embedded text
                text = page.get_text("text") or ""
                text = (text or "").strip()

                # 2) OCR only if no embedded text
                if not text and USE_OCR:
                    pix = page.get_pixmap(dpi=300)  # higher DPI for slides
                    img = Image.open(io.BytesIO(pix.tobytes("png"))).convert(
                        "L"
                    )  # grayscale often helps OCR
                    try:
                        text = pytesseract.image_to_string(
                            img, lang="eng", config="--oem 1 --psm 6"
                        )
                    except Exception:
                        text = ""

                text = (text or "").strip()
                if not text:
                    continue  # skip empty chunks

                locator = f"page-{i}"
                self.chunks[f"{base_uri}#{locator}"] = text
        finally:
            doc.close()

    def _parse_pptx(self, path: Path, base_uri: str):
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            for sp in slide.shapes:
                # text frames
                if hasattr(sp, "text") and sp.text:
                    t = (sp.text or "").strip()
                    if t:
                        texts.append(t)
                # tables
                if hasattr(sp, "has_table") and sp.has_table:
                    try:
                        for row in sp.table.rows:
                            cells = [(cell.text or "").strip() for cell in row.cells]
                            line = " | ".join([c for c in cells if c])
                            if line:
                                texts.append(line)
                    except Exception:
                        pass

            # notes
            notes = ""
            try:
                if (
                    slide.has_notes_slide
                    and slide.notes_slide
                    and slide.notes_slide.notes_text_frame
                ):
                    notes = slide.notes_slide.notes_text_frame.text or ""
                    notes = (notes or "").strip()
            except Exception:
                notes = ""

            merged = "\n".join(texts + ([f"Notes:\n{notes}"] if notes else []))
            merged = (merged or "").strip()
            if not merged:
                continue

            locator = f"slide-{i}"
            self.chunks[f"{base_uri}#{locator}"] = merged

    def _parse_docx(self, path: Path, base_uri: str):
        doc = Document(path)

        # headers & footers text collector
        def _collect_section_parts(d):
            out = []
            for s in d.sections:
                try:
                    if s.header:
                        out.extend(
                            [
                                (p.text or "").strip()
                                for p in s.header.paragraphs
                                if (p.text or "").strip()
                            ]
                        )
                except Exception:
                    pass
                try:
                    if s.footer:
                        out.extend(
                            [
                                (p.text or "").strip()
                                for p in s.footer.paragraphs
                                if (p.text or "").strip()
                            ]
                        )
                except Exception:
                    pass
            return "\n".join(out).strip()

        header_footer = _collect_section_parts(doc)

        sec_idx = 0
        buf: List[str] = []
        title = "section"
        ### FIX 2: Add counter for smart-splitting on blank lines ###
        blank_line_count = 0

        def save_buffer(current_buf, current_title, current_idx):
            """Helper to save the current buffer as a chunk."""
            chunk_text = "\n".join(current_buf).strip()
            # Filter out tiny/useless chunks
            if chunk_text and (len(chunk_text) > 80 or '\n' in chunk_text):
                self.chunks[f"{base_uri}#{current_title}-{current_idx}"] = chunk_text
            return [], current_idx + 1 # Return new buffer and incremented index

        # paragraphs (with simple heading-based segmentation)
        for p in doc.paragraphs:
            txt = (p.text or "").strip()
            
            is_heading = p.style and getattr(p.style, "name", "").startswith("Heading")
            is_blank = not txt

            # --- Check for chunk-splitting conditions ---
            # Condition 1: We hit a heading
            # Condition 2: We hit two consecutive blank lines (split on the 2nd line)
            should_split = (
                is_heading or
                (is_blank and blank_line_count >= 1)
            )

            if should_split and buf:
                # Save the previous buffer as a chunk
                buf, sec_idx = save_buffer(buf, title, sec_idx)
                blank_line_count = 0 # Reset counter

                if is_heading:
                    title = txt or "section" # New title from heading
                else:
                    title = "section" # Reset to default title
            
            # --- Handle the current paragraph ---
            if is_blank:
                blank_line_count += 1
            else:
                blank_line_count = 0 # Reset on any non-blank line
                if not is_heading: # Don't add the heading itself to the buffer
                    buf.append(txt)

        # tables → append their rows as lines
        for table in doc.tables:
            for row in table.rows:
                cells = [(cell.text or "").strip() for cell in row.cells]
                line = " | ".join([c for c in cells if c])
                if line:
                    buf.append(line)

        if header_footer:
            buf.append("\n" + header_footer)

        # Save any remaining text in the buffer
        if buf:
            buf, sec_idx = save_buffer(buf, title, sec_idx)

    def _parse_text(self, path: Path, base_uri: str):
        try:
            txt = path.read_text(encoding="utf-8", errors="ignore")
        except Exception:
            txt = path.read_text(errors="ignore")
        txt = (txt or "").strip()
        if not txt:
            return
        self.chunks[f"{base_uri}#full"] = txt

    def _parse_image(self, path: Path, base_uri: str):
        text = ""
        if USE_OCR:
            try:
                img = Image.open(path).convert("L")
                text = pytesseract.image_to_string(
                    img, lang="eng", config="--oem 1 --psm 6"
                )
            except Exception:
                text = ""
        text = (text or "").strip()
        if not text:
            return
        self.chunks[f"{base_uri}#image"] = text


# ------------ build index at import ------------
INDEX = Index()
INDEX.index_all()